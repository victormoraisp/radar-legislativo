"""
Gera a apresentação executiva do Radar Legislativo (6 slides, formato pitch).

Uso:
    python apresentacao/build_pptx.py

Saída:
    apresentacao/radar_legislativo_pitch.pptx

Se os prints existirem em docs/img/, eles são inseridos automaticamente:
    docs/img/print_supabase_tabelas.png  (slide 4)
    docs/img/print_n8n_execucao.png      (slide 5)
    docs/img/print_email_resumo.png      (slide 5)
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt
from pptx.enum.dml import MSO_LINE_DASH_STYLE


BASE_DIR = Path(__file__).resolve().parent.parent
IMG_DIR = BASE_DIR / "docs" / "img"
SAIDA = Path(__file__).resolve().parent / "radar_legislativo_pitch.pptx"

# Paleta
AZUL_ESCURO = RGBColor(0x0F, 0x17, 0x2A)
AZUL = RGBColor(0x1D, 0x4E, 0xD8)
AZUL_CLARO = RGBColor(0xDB, 0xEA, 0xFE)
VERDE = RGBColor(0x15, 0x80, 0x3D)
LARANJA = RGBColor(0xEA, 0x58, 0x0C)
ROXO = RGBColor(0x7C, 0x3A, 0xED)
CINZA = RGBColor(0x47, 0x55, 0x69)
CINZA_CLARO = RGBColor(0xF1, 0xF5, 0xF9)
BRANCO = RGBColor(0xFF, 0xFF, 0xFF)

LARGURA = Inches(13.333)
ALTURA = Inches(7.5)


def novo_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # layout em branco


def fundo(slide, cor):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = cor


def caixa_texto(slide, x, y, largura, altura):
    caixa = slide.shapes.add_textbox(x, y, largura, altura)
    caixa.text_frame.word_wrap = True
    return caixa


def paragrafo(tf, texto, tamanho, cor, negrito=False, primeiro=False,
              alinhamento=PP_ALIGN.LEFT, espaco_antes=6):
    p = tf.paragraphs[0] if primeiro else tf.add_paragraph()
    p.text = texto
    p.alignment = alinhamento
    p.space_before = Pt(espaco_antes)
    p.font.size = Pt(tamanho)
    p.font.color.rgb = cor
    p.font.bold = negrito
    return p


def titulo_slide(slide, texto, subtitulo=None):
    faixa = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, LARGURA, Inches(1.1))
    faixa.fill.solid()
    faixa.fill.fore_color.rgb = AZUL_ESCURO
    faixa.line.fill.background()

    tf = faixa.text_frame
    tf.margin_left = Inches(0.5)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragrafo(tf, texto, 28, BRANCO, negrito=True, primeiro=True)

    if subtitulo:
        paragrafo(tf, subtitulo, 13, RGBColor(0xCB, 0xD5, 0xE1))


def bloco(slide, x, y, largura, altura, cor_fundo, cor_borda,
          titulo=None, linhas=None, tamanho_titulo=15, tamanho_linha=12):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, largura, altura)
    shape.fill.solid()
    shape.fill.fore_color.rgb = cor_fundo
    shape.line.color.rgb = cor_borda
    shape.line.width = Pt(1.5)

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    if titulo:
        paragrafo(tf, titulo, tamanho_titulo, AZUL_ESCURO, negrito=True,
                  primeiro=True, alinhamento=PP_ALIGN.CENTER, espaco_antes=0)

    for linha in (linhas or []):
        paragrafo(tf, linha, tamanho_linha, CINZA,
                  alinhamento=PP_ALIGN.CENTER, espaco_antes=2)

    return shape


def seta(slide, x, y, largura=Inches(0.35)):
    forma = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, y, largura, Inches(0.3))
    forma.fill.solid()
    forma.fill.fore_color.rgb = CINZA
    forma.line.fill.background()
    return forma


def imagem_ou_placeholder(slide, nome_arquivo, x, y, largura, altura, legenda):
    caminho = IMG_DIR / nome_arquivo

    if caminho.exists():
        slide.shapes.add_picture(str(caminho), x, y, width=largura)
    else:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, largura, altura)
        shape.fill.solid()
        shape.fill.fore_color.rgb = CINZA_CLARO
        shape.line.color.rgb = CINZA
        shape.line.dash_style = MSO_LINE_DASH_STYLE.DASH
        tf = shape.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        paragrafo(tf, f"[Inserir print: {legenda}]", 12, CINZA,
                  primeiro=True, alinhamento=PP_ALIGN.CENTER)

    caixa = caixa_texto(slide, x, y + altura + Inches(0.05), largura, Inches(0.3))
    paragrafo(caixa.text_frame, legenda, 10, CINZA, primeiro=True,
              alinhamento=PP_ALIGN.CENTER, espaco_antes=0)


def slide_1(prs):
    slide = novo_slide(prs)
    fundo(slide, AZUL_ESCURO)

    caixa = caixa_texto(slide, Inches(0.8), Inches(0.7), Inches(11.7), Inches(1.8))
    tf = caixa.text_frame
    paragrafo(tf, "RADAR LEGISLATIVO", 44, BRANCO, negrito=True, primeiro=True)
    paragrafo(tf, "Inteligência legislativa automatizada para a Bússola Pública", 20,
              RGBColor(0x93, 0xC5, 0xFD))

    caixa = caixa_texto(slide, Inches(0.8), Inches(2.6), Inches(11.7), Inches(0.5))
    paragrafo(caixa.text_frame, "O PROBLEMA", 18, LARANJA, negrito=True, primeiro=True)

    problemas = [
        ("2 analistas, 0 escala", "Leitura manual do site da Câmara o dia inteiro para montar um relatório vendido a R$ 15 mil/mês por cliente."),
        ("Nenhuma base de dados", "Planilhas pessoais, sem histórico organizado. Tudo o que foi consultado se perde."),
        ("Classificação inconsistente", "Cada analista rotula como quer. Nenhum indicador é medido."),
        ("Alertas por memória", "Se o analista esquece, o cliente é pego de surpresa por projetos que tramitavam há meses."),
    ]

    for i, (titulo, texto) in enumerate(problemas):
        x = Inches(0.8 + (i % 2) * 6.0)
        y = Inches(3.2 + (i // 2) * 1.85)
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(5.7), Inches(1.65))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0x1E, 0x29, 0x3B)
        shape.line.color.rgb = RGBColor(0x33, 0x41, 0x55)
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.2)
        tf.margin_right = Inches(0.2)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        paragrafo(tf, titulo, 15, RGBColor(0xFD, 0xBA, 0x74), negrito=True, primeiro=True, espaco_antes=0)
        paragrafo(tf, texto, 12, RGBColor(0xCB, 0xD5, 0xE1), espaco_antes=3)

    caixa = caixa_texto(slide, Inches(0.8), Inches(6.95), Inches(11.7), Inches(0.4))
    paragrafo(caixa.text_frame,
              "O problema não é falta de dado. É falta de engenharia de dados.",
              14, BRANCO, negrito=True, primeiro=True, alinhamento=PP_ALIGN.CENTER)

    return slide


def slide_2(prs):
    slide = novo_slide(prs)
    fundo(slide, BRANCO)
    titulo_slide(slide, "A solução: pipeline automatizado com IA",
                 "Do dado público bruto ao relatório executivo, sem intervenção manual")

    entregas = [
        ("Base única e histórica", "PostgreSQL na nuvem (Supabase) com 30 dias de dados: 16.082 proposições, 1.393 votações, despesas e cadastros.", AZUL, AZUL_CLARO),
        ("IA que enriquece o dado", "Resumo executivo por proposição gerado pela OpenAI e persistido no banco (coluna resumo_ia). Não é decoração: aparece no relatório do cliente.", ROXO, RGBColor(0xED, 0xE9, 0xFE)),
        ("Entrega automática às 06h", "Workflow n8n consulta o banco, gera análise executiva com IA e envia o resumo diário por e-mail. Roda sozinho, todo dia.", LARANJA, RGBColor(0xFF, 0xED, 0xD5)),
        ("Rastreável e reprocessável", "JSON bruto preservado, validações com Pandas, versionamento no GitHub e custo de IA medido antes de escalar.", VERDE, RGBColor(0xDC, 0xFC, 0xE7)),
    ]

    for i, (titulo, texto, cor, cor_fundo) in enumerate(entregas):
        x = Inches(0.7 + (i % 2) * 6.1)
        y = Inches(1.5 + (i // 2) * 2.5)
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(5.9), Inches(2.25))
        shape.fill.solid()
        shape.fill.fore_color.rgb = cor_fundo
        shape.line.color.rgb = cor
        shape.line.width = Pt(1.75)
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.25)
        tf.margin_right = Inches(0.25)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        paragrafo(tf, titulo, 17, AZUL_ESCURO, negrito=True, primeiro=True, espaco_antes=0)
        paragrafo(tf, texto, 13, CINZA, espaco_antes=5)

    caixa = caixa_texto(slide, Inches(0.7), Inches(6.7), Inches(12), Inches(0.5))
    paragrafo(caixa.text_frame,
              "Dois analistas deixam de ler site e passam a validar sinal. O produto de R$ 15 mil/mês escala sem contratar.",
              14, AZUL, negrito=True, primeiro=True, alinhamento=PP_ALIGN.CENTER)

    return slide


def slide_3(prs):
    slide = novo_slide(prs)
    fundo(slide, BRANCO)
    titulo_slide(slide, "Arquitetura e stack",
                 "Python · Pandas · SQL · Supabase · OpenAI · n8n · Git/GitHub")

    y1 = Inches(1.6)
    etapas_linha1 = [
        ("API Câmara", ["Dados Abertos", "5 endpoints"], RGBColor(0x02, 0x84, 0xC7), RGBColor(0xE0, 0xF2, 0xFE)),
        ("Extração Python", ["requests + paginação", "try/except + retry"], AZUL, AZUL_CLARO),
        ("JSON bruto", ["data/raw", "reprocessável"], RGBColor(0xCA, 0x8A, 0x04), RGBColor(0xFE, 0xF9, 0xC3)),
        ("Pandas", ["validações, tipagem", "deduplicação"], VERDE, RGBColor(0xDC, 0xFC, 0xE7)),
    ]

    x = Inches(0.55)
    for i, (titulo, linhas, cor, cor_fundo) in enumerate(etapas_linha1):
        bloco(slide, x, y1, Inches(2.65), Inches(1.35), cor_fundo, cor, titulo, linhas)
        x += Inches(2.65)
        if i < 3:
            seta(slide, x + Inches(0.03), y1 + Inches(0.52), Inches(0.3))
            x += Inches(0.38)

    seta_baixo = slide.shapes.add_shape(
        MSO_SHAPE.DOWN_ARROW, Inches(11.1), Inches(3.05), Inches(0.3), Inches(0.4))
    seta_baixo.fill.solid()
    seta_baixo.fill.fore_color.rgb = CINZA
    seta_baixo.line.fill.background()

    y2 = Inches(3.55)
    bloco(slide, Inches(7.3), y2, Inches(5.4), Inches(1.5),
          RGBColor(0xED, 0xE9, 0xFE), ROXO,
          "Supabase · PostgreSQL (schema radar)",
          ["3 tabelas fato + 4 dimensões", "views SQL analíticas", "coluna resumo_ia gerada por IA"])

    bloco(slide, Inches(3.3), y2, Inches(3.4), Inches(1.5),
          RGBColor(0xFC, 0xE7, 0xF3), RGBColor(0xDB, 0x27, 0x77),
          "Camada de IA (OpenAI)",
          ["resumo executivo", "por proposição", "custo medido antes de escalar"])

    seta(slide, Inches(6.75), y2 + Inches(0.6), Inches(0.5))

    seta_baixo2 = slide.shapes.add_shape(
        MSO_SHAPE.DOWN_ARROW, Inches(9.85), Inches(5.15), Inches(0.3), Inches(0.4))
    seta_baixo2.fill.solid()
    seta_baixo2.fill.fore_color.rgb = CINZA
    seta_baixo2.line.fill.background()

    y3 = Inches(5.65)
    bloco(slide, Inches(6.4), y3, Inches(6.2), Inches(1.4),
          RGBColor(0xFF, 0xED, 0xD5), LARANJA,
          "n8n · automação diária às 06h",
          ["Supabase REST → contexto → OpenAI → HTML → e-mail", "workflow exportado e versionado no GitHub"])

    bloco(slide, Inches(3.3), y3, Inches(2.5), Inches(1.4),
          CINZA_CLARO, CINZA,
          "Cliente",
          ["resumo diário", "no e-mail"])

    seta_esq = slide.shapes.add_shape(
        MSO_SHAPE.LEFT_ARROW, Inches(5.85), y3 + Inches(0.55), Inches(0.5), Inches(0.3))
    seta_esq.fill.solid()
    seta_esq.fill.fore_color.rgb = CINZA
    seta_esq.line.fill.background()

    return slide


def slide_4(prs):
    slide = novo_slide(prs)
    fundo(slide, BRANCO)
    titulo_slide(slide, "Demo · Banco populado e modelo de dados",
                 "Supabase/PostgreSQL com 30 dias de dados ingeridos (abril/2026)")

    dados = [
        ("fato_proposicao", "16.082", "com resumo_ia"),
        ("fato_votacao", "1.393", ""),
        ("fato_despesa", "3.000", "amostra de 36.542"),
        ("dim_deputado", "512", ""),
        ("dim_partido", "21", ""),
        ("dim_tema + dim_situacao", "131", "referências"),
    ]

    for i, (tabela, linhas, obs) in enumerate(dados):
        x = Inches(0.6)
        y = Inches(1.55 + i * 0.92)
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(5.2), Inches(0.8))
        shape.fill.solid()
        shape.fill.fore_color.rgb = CINZA_CLARO
        shape.line.color.rgb = AZUL
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.2)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = paragrafo(tf, f"{tabela}:  {linhas} linhas", 14, AZUL_ESCURO, negrito=True,
                      primeiro=True, espaco_antes=0)
        if obs:
            run = p.add_run()
            run.text = f"   ({obs})"
            run.font.size = Pt(11)
            run.font.color.rgb = CINZA
            run.font.bold = False

    imagem_ou_placeholder(
        slide, "print_supabase_tabelas.png",
        Inches(6.2), Inches(1.55), Inches(6.5), Inches(4.6),
        "Table Editor do Supabase — schema radar populado")

    caixa = caixa_texto(slide, Inches(0.6), Inches(6.9), Inches(12.1), Inches(0.45))
    paragrafo(caixa.text_frame,
              "Modelo estrela: 3 fatos + 4 dimensões · validações com Pandas (nulos, datas, valores) · acesso em modo leitura no README",
              12, CINZA, primeiro=True, alinhamento=PP_ALIGN.CENTER)

    return slide


def slide_5(prs):
    slide = novo_slide(prs)
    fundo(slide, BRANCO)
    titulo_slide(slide, "Demo · Automação n8n + IA no produto",
                 "A IA enriquece o dado no banco e aparece no relatório do cliente")

    imagem_ou_placeholder(
        slide, "print_n8n_execucao.png",
        Inches(0.6), Inches(1.55), Inches(6.0), Inches(3.3),
        "Execução bem-sucedida do workflow n8n (diário às 06h)")

    imagem_ou_placeholder(
        slide, "print_email_resumo.png",
        Inches(6.9), Inches(1.55), Inches(5.8), Inches(3.3),
        "E-mail recebido: análise executiva + coluna Resumo IA")

    caixa = caixa_texto(slide, Inches(0.6), Inches(5.4), Inches(12.1), Inches(1.7))
    tf = caixa.text_frame
    paragrafo(tf, "IA aplicada em dois pontos:", 15, AZUL_ESCURO, negrito=True, primeiro=True)
    paragrafo(tf, "1. Resumo executivo por proposição persistido na coluna resumo_ia (Caminho B do desafio) — exibido na tabela do e-mail.", 13, CINZA)
    paragrafo(tf, "2. Análise executiva do dia gerada no n8n a partir das últimas proposições e votações — abertura do e-mail.", 13, CINZA)
    paragrafo(tf, "Custo sob controle: teste com 10 registros, custo real impresso e projeção antes de escalar para as 16 mil proposições.", 13, VERDE, negrito=True)

    return slide


def slide_6(prs):
    slide = novo_slide(prs)
    fundo(slide, AZUL_ESCURO)

    caixa = caixa_texto(slide, Inches(0.8), Inches(0.6), Inches(11.7), Inches(0.8))
    paragrafo(caixa.text_frame, "Próximos passos", 32, BRANCO, negrito=True, primeiro=True)

    passos = [
        ("Ingestão incremental diária", "O n8n passa a acionar também a carga: dado novo todo dia, histórico crescente."),
        ("Classificação temática com embeddings", "pgvector já habilitado no Supabase: tema automático por similaridade de cosseno (Caminho A)."),
        ("Alertas por tema crítico", "Telegram/e-mail imediato quando proposição de Tecnologia, Tributário ou setor do cliente é apresentada."),
        ("Dashboard executivo", "Power BI conectado ao Supabase: proposições por tema, atividade por partido e deputado."),
        ("Produto multi-cliente", "Relatórios segmentados por setor de interesse: o relatório de R$ 15 mil/mês vira linha de produção."),
    ]

    for i, (titulo, texto) in enumerate(passos):
        y = Inches(1.6 + i * 0.98)
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), y, Inches(11.7), Inches(0.85))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0x1E, 0x29, 0x3B)
        shape.line.color.rgb = RGBColor(0x33, 0x41, 0x55)
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.25)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = paragrafo(tf, f"{i + 1}. {titulo}", 14, RGBColor(0x93, 0xC5, 0xFD), negrito=True,
                      primeiro=True, espaco_antes=0)
        run = p.add_run()
        run.text = f"  —  {texto}"
        run.font.size = Pt(12)
        run.font.color.rgb = RGBColor(0xCB, 0xD5, 0xE1)
        run.font.bold = False

    caixa = caixa_texto(slide, Inches(0.8), Inches(6.75), Inches(11.7), Inches(0.5))
    paragrafo(caixa.text_frame,
              "Radar Legislativo · github.com/victormoraisp/radar-legislativo · Victor Morais",
              13, BRANCO, primeiro=True, alinhamento=PP_ALIGN.CENTER)

    return slide


def main():
    prs = Presentation()
    prs.slide_width = LARGURA
    prs.slide_height = ALTURA

    slide_1(prs)
    slide_2(prs)
    slide_3(prs)
    slide_4(prs)
    slide_5(prs)
    slide_6(prs)

    prs.save(SAIDA)
    print(f"Apresentação gerada: {SAIDA}")

    faltando = [
        nome for nome in [
            "print_supabase_tabelas.png",
            "print_n8n_execucao.png",
            "print_email_resumo.png",
        ] if not (IMG_DIR / nome).exists()
    ]

    if faltando:
        print("\nPrints ausentes (placeholders foram usados no lugar):")
        for nome in faltando:
            print(f"  - docs/img/{nome}")
        print("Salve os prints e rode o script novamente para incorporá-los.")


if __name__ == "__main__":
    main()
